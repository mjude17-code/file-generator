from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from io import BytesIO
from datetime import datetime

from docx import Document
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import openpyxl
from openpyxl.styles import Font
from pptx import Presentation
from pptx.util import Inches

app = Flask(__name__)
CORS(app)

@app.route('/')
def home():
    return jsonify({'message': 'File Generator API', 'ui': '/ui'})

@app.route('/ui')
def ui():
    return '''
<!DOCTYPE html>
<html>
<head>
    <title>File Generator</title>
    <style>
        body { font-family: Arial; max-width: 800px; margin: 50px auto; background: #f5f5f5; }
        .container { background: white; padding: 40px; border-radius: 10px; box-shadow: 0 0 10px rgba(0,0,0,0.1); }
        h1 { color: #333; }
        .form-group { margin: 20px 0; }
        label { font-weight: bold; display: block; margin-bottom: 5px; }
        input, textarea, select { width: 100%; padding: 10px; margin-bottom: 10px; border: 1px solid #ddd; border-radius: 5px; font-size: 14px; }
        button { width: 100%; padding: 12px; background: #007bff; color: white; border: none; border-radius: 5px; cursor: pointer; font-size: 16px; font-weight: bold; }
        button:hover { background: #0056b3; }
        .result { margin-top: 20px; padding: 15px; border-radius: 5px; }
        .success { background: #d4edda; color: #155724; }
        .error { background: #f8d7da; color: #721c24; }
        .hidden { display: none; }
    </style>
</head>
<body>
    <div class="container">
        <h1>📄 File Generator</h1>
        
        <div class="form-group">
            <label>File Type:</label>
            <select id="fileType" onchange="toggleSections()">
                <option value="">-- Select File Type --</option>
                <option value="word">Word (.docx)</option>
                <option value="pdf">PDF</option>
                <option value="excel">Excel (.xlsx)</option>
                <option value="ppt">PowerPoint (.pptx)</option>
            </select>
        </div>
        
        <div class="form-group">
            <label>Title:</label>
            <input type="text" id="title" placeholder="Enter document title" />
        </div>
        
        <div class="form-group">
            <label>Content:</label>
            <textarea id="content" placeholder="Paste your content here" rows="8"></textarea>
        </div>
        
        <div id="excelSection" class="form-group hidden">
            <label>Excel Rows (JSON):</label>
            <textarea id="rows" placeholder='[["Name","Age"],["John",30],["Jane",25]]' rows="6"></textarea>
        </div>
        
        <div id="pptSection" class="form-group hidden">
            <label>PowerPoint Slides (JSON):</label>
            <textarea id="slides" placeholder='[{"heading":"Slide 1","content":"Welcome"}]' rows="6"></textarea>
        </div>
        
        <button onclick="generateFile()">🚀 Generate & Download File</button>
        
        <div id="result"></div>
    </div>
    
    <script>
        function toggleSections() {
            const fileType = document.getElementById('fileType').value;
            document.getElementById('excelSection').classList.toggle('hidden', fileType !== 'excel');
            document.getElementById('pptSection').classList.toggle('hidden', fileType !== 'ppt');
        }
        
        function generateFile() {
            const fileType = document.getElementById('fileType').value;
            const title = document.getElementById('title').value;
            const content = document.getElementById('content').value;
            const resultDiv = document.getElementById('result');
            
            if (!fileType) {
                resultDiv.innerHTML = '<div class="result error">❌ Please select file type</div>';
                return;
            }
            if (!title) {
                resultDiv.innerHTML = '<div class="result error">❌ Please enter title</div>';
                return;
            }
            
            let payload = { file_type: fileType, title, content };
            
            try {
                if (fileType === 'excel') {
                    payload.rows = JSON.parse(document.getElementById('rows').value || '[]');
                }
                if (fileType === 'ppt') {
                    payload.slides = JSON.parse(document.getElementById('slides').value || '[]');
                }
            } catch (e) {
                resultDiv.innerHTML = '<div class="result error">❌ Invalid JSON: ' + e.message + '</div>';
                return;
            }
            
            resultDiv.innerHTML = '<div class="result">⏳ Generating...</div>';
            
            fetch('/generate', {
                method: 'POST',
                headers: { 'Content-Type': 'application/json' },
                body: JSON.stringify(payload)
            })
            .then(r => r.blob())
            .then(blob => {
                const url = window.URL.createObjectURL(blob);
                const a = document.createElement('a');
                a.href = url;
                const ext = {word: '.docx', pdf: '.pdf', excel: '.xlsx', ppt: '.pptx'}[fileType];
                a.download = title.replace(/ /g, '_') + ext;
                a.click();
                window.URL.revokeObjectURL(url);
                resultDiv.innerHTML = '<div class="result success">✅ File downloaded!</div>';
            })
            .catch(e => {
                resultDiv.innerHTML = '<div class="result error">❌ Error: ' + e.message + '</div>';
            });
        }
    </script>
</body>
</html>
    '''

@app.route('/generate', methods=['POST'])
def generate_universal():
    try:
        data = request.json
        file_type = data.get('file_type', 'word').lower()
        
        if file_type == 'word':
            return generate_word(data)
        elif file_type == 'pdf':
            return generate_pdf(data)
        elif file_type == 'excel':
            return generate_excel(data)
        elif file_type == 'ppt':
            return generate_ppt(data)
        else:
            return jsonify({'error': 'Invalid file_type'}), 400
    except Exception as e:
        return jsonify({'error': str(e)}), 500

def generate_word(data):
    title = data.get('title', 'Document')
    content = data.get('content', '')
    doc = Document()
    doc.add_heading(title, 0)
    doc.add_paragraph(content)
    file_stream = BytesIO()
    doc.save(file_stream)
    file_stream.seek(0)
    return send_file(file_stream, mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document', as_attachment=True, download_name=f'{title.replace(" ", "_")}.docx')

def generate_pdf(data):
    title = data.get('title', 'Document')
    content = data.get('content', '')
    file_stream = BytesIO()
    pdf_canvas = canvas.Canvas(file_stream, pagesize=letter)
    pdf_canvas.setFont("Helvetica-Bold", 16)
    pdf_canvas.drawString(50, 750, title)
    pdf_canvas.setFont("Helvetica", 12)
    y_position = 720
    for line in content.split('\n'):
        pdf_canvas.drawString(50, y_position, line)
        y_position -= 20
    pdf_canvas.save()
    file_stream.seek(0)
    return send_file(file_stream, mimetype='application/pdf', as_attachment=True, download_name=f'{title.replace(" ", "_")}.pdf')

def generate_excel(data):
    title = data.get('title', 'Sheet')
    rows = data.get('rows', [])
    wb = openpyxl.Workbook()
    ws = wb.active
    ws['A1'] = title
    ws['A1'].font = Font(bold=True, size=14)
    for idx, row in enumerate(rows, start=2):
        for col_idx, value in enumerate(row, start=1):
            ws.cell(row=idx, column=col_idx, value=value)
    file_stream = BytesIO()
    wb.save(file_stream)
    file_stream.seek(0)
    return send_file(file_stream, mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', as_attachment=True, download_name=f'{title.replace(" ", "_")}.xlsx')

def generate_ppt(data):
    title = data.get('title', 'Presentation')
    slides_data = data.get('slides', [])
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = f"Generated on {datetime.now().strftime('%Y-%m-%d')}"
    for slide_data in slides_data:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        slide.shapes.title.text = slide_data.get('heading', 'Slide')
        slide.placeholders[1].text_frame.text = slide_data.get('content', '')
    file_stream = BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)
    return send_file(file_stream, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation', as_attachment=True, download_name=f'{title.replace(" ", "_")}.pptx')

@app.route('/health', methods=['GET'])
def health():
    return jsonify({'status': 'ok'}), 200

if __name__ == '__main__':
    app.run(debug=True, port=8000)
