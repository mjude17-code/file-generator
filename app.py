from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
import os
from io import BytesIO
from datetime import datetime

# File generation libraries
from docx import Document
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import openpyxl
from openpyxl.styles import Font, PatternFill
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

app = Flask(__name__)
CORS(app)

# ===== HOME ROUTE =====
@app.route("/")
def home():
    return jsonify({
        "message": "File Generator API",
        "endpoints": {
            "word": "/generate/word",
            "pdf": "/generate/pdf",
            "excel": "/generate/excel",
            "ppt": "/generate/ppt",
            "health": "/health"
        }
    })

# ===== WORD FILE GENERATION =====
@app.route('/generate/word', methods=['POST'])
def generate_word():
    try:
        data = request.json
        title = data.get('title', 'Document')
        content = data.get('content', 'Default content')
        
        doc = Document()
        doc.add_heading(title, 0)
        doc.add_paragraph(content)
        
        # Save to BytesIO (in-memory file)
        file_stream = BytesIO()
        doc.save(file_stream)
        file_stream.seek(0)
        
        return send_file(
            file_stream,
            mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            as_attachment=True,
            download_name=f'{title.replace(" ", "_")}.docx'
        )
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ===== PDF FILE GENERATION =====
@app.route('/generate/pdf', methods=['POST'])
def generate_pdf():
    try:
        data = request.json
        title = data.get('title', 'Document')
        content = data.get('content', 'Default content')
        
        file_stream = BytesIO()
        pdf_canvas = canvas.Canvas(file_stream, pagesize=letter)
        
        # Title
        pdf_canvas.setFont("Helvetica-Bold", 16)
        pdf_canvas.drawString(50, 750, title)
        
        # Content
        pdf_canvas.setFont("Helvetica", 12)
        y_position = 720
        for line in content.split('\n'):
            pdf_canvas.drawString(50, y_position, line)
            y_position -= 20
        
        pdf_canvas.save()
        file_stream.seek(0)
        
        return send_file(
            file_stream,
            mimetype='application/pdf',
            as_attachment=True,
            download_name=f'{title.replace(" ", "_")}.pdf'
        )
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ===== EXCEL FILE GENERATION =====
@app.route('/generate/excel', methods=['POST'])
def generate_excel():
    try:
        data = request.json
        title = data.get('title', 'Sheet')
        rows = data.get('rows', [])  # List of lists
        
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = 'Sheet1'
        
        # Add title
        ws['A1'] = title
        ws['A1'].font = Font(bold=True, size=14)
        
        # Add rows
        for idx, row in enumerate(rows, start=2):
            for col_idx, value in enumerate(row, start=1):
                ws.cell(row=idx, column=col_idx, value=value)
        
        file_stream = BytesIO()
        wb.save(file_stream)
        file_stream.seek(0)
        
        return send_file(
            file_stream,
            mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            as_attachment=True,
            download_name=f'{title.replace(" ", "_")}.xlsx'
        )
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ===== POWERPOINT FILE GENERATION =====
@app.route('/generate/ppt', methods=['POST'])
def generate_ppt():
    try:
        data = request.json
        title = data.get('title', 'Presentation')
        slides_data = data.get('slides', [])  # List of dicts with 'heading' and 'content'
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        title_shape.text = title
        subtitle_shape.text = f"Generated on {datetime.now().strftime('%Y-%m-%d')}"
        
        # Content slides
        for slide_data in slides_data:
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            title_shape = slide.shapes.title
            body_shape = slide.placeholders[1]
            
            title_shape.text = slide_data.get('heading', 'Slide')
            text_frame = body_shape.text_frame
            text_frame.text = slide_data.get('content', '')
        
        file_stream = BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        
        return send_file(
            file_stream,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=f'{title.replace(" ", "_")}.pptx'
        )
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ===== HEALTH CHECK =====
@app.route('/health', methods=['GET'])
def health():
    return jsonify({'status': 'ok'}), 200

if __name__ == '__main__':
    app.run(debug=True, port=8000)
